"""Shared PPTX builder utilities."""

from __future__ import annotations

import os
from typing import Sequence

from pptx import Presentation
from pptx.util import Cm, Pt

from my_chart.config import PPTX_HEIGHT, PPTX_WIDTH


# @MX:ANCHOR: [AUTO] PPTX factory - fan_in=3, used by queries, momentum, market analysis
# @MX:REASON: All PPTX exports start here; changing dimensions affects every generated report
def create_widescreen_pptx() -> Presentation:
    """Create a 16:9 widescreen PPTX presentation."""
    prs = Presentation()
    prs.slide_width = PPTX_WIDTH
    prs.slide_height = PPTX_HEIGHT
    return prs


def add_image_slide(
    prs: Presentation,
    image_path: str,
    links: dict[str, str] | None = None,
) -> None:
    """Add a full-size image slide with optional hyperlinks."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(image_path, 0, 0, PPTX_WIDTH, PPTX_HEIGHT)

    if links:
        y_offset = Cm(0.5)
        for label, url in links.items():
            txBox = slide.shapes.add_textbox(Cm(0), y_offset, Cm(4), Cm(0.6))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = label
            r.hyperlink.address = url
            r.font.size = Pt(10)
            y_offset += Cm(0.6)


def add_grid_slide(
    prs: Presentation,
    image_paths: Sequence[str],
    links_per_image: Sequence[dict[str, str]] | None = None,
) -> None:
    """Add a 2x2 grid slide with images and optional links."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    pic_w = PPTX_WIDTH / 2
    pic_h = PPTX_HEIGHT / 2

    for j, file in enumerate(image_paths):
        row, col = j // 2, j % 2
        top = pic_h * row
        left = pic_w * col
        slide.shapes.add_picture(file, left, top, pic_w, pic_h)

        if links_per_image and j < len(links_per_image):
            y_offset = top
            for label, url in links_per_image[j].items():
                txBox = slide.shapes.add_textbox(left, y_offset, Cm(3), Cm(0.6))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = label
                r.hyperlink.address = url
                r.font.size = Pt(10)
                y_offset += Cm(0.6)


def save_and_cleanup(
    prs: Presentation,
    output_path: str,
    temp_files: Sequence[str],
) -> None:
    """Save PPTX and remove temporary image files."""
    prs.save(output_path)
    for f in temp_files:
        try:
            os.remove(f)
        except FileNotFoundError:
            pass
